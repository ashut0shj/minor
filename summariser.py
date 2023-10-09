# -*- coding: utf-8 -*-
"""minor project.ipynb

Automatically generated by Colaboratory.

Original file is located at
    https://colab.research.google.com/drive/1jvIblukSGj12nU8h1R-2eVeULLRxKTRv
"""

import spacy
from spacy.lang.en.stop_words import STOP_WORDS
from string import punctuation
from heapq import nlargest

stopwords = list(STOP_WORDS)
nlp = spacy.load('en_core_web_sm')

text = input("Enter the text you want to summarize:\n")

doc = nlp(text)
tokens = [token.text for token in doc]

punctuation = punctuation + '\n'
word_frequencies = {}

for word in doc:
    if word.text.lower() not in stopwords and word.text.lower() not in punctuation:
        if word.text not in word_frequencies.keys():
            word_frequencies[word.text] = 1
        else:
            word_frequencies[word.text] += 1

max_frequency = max(word_frequencies.values())

for word in word_frequencies.keys():
    word_frequencies[word] = word_frequencies[word] / max_frequency

sentence_tokens = [sent for sent in doc.sents]
sentence_scores = {}

for sent in sentence_tokens:
    for word in sent:
        if word.text.lower() in word_frequencies.keys():
            if sent not in sentence_scores.keys():
                sentence_scores[sent] = word_frequencies[word.text.lower()]
            else:
                sentence_scores[sent] += word_frequencies[word.text.lower()]

select_length = int(len(sentence_tokens) * 0.3)
summary = nlargest(select_length, sentence_scores, key=sentence_scores.get)
final_summary = [sent.text for sent in summary]

summary = ' '.join(final_summary)
print("\nSummary:\n", summary)

import pandas as pd
import numpy as np
import re
from pptx import Presentation
from io import BytesIO
from google.colab import files

def process_pptx_content(pptx_content):
    text_runs = []

    with BytesIO(pptx_content) as f:
        prs = Presentation(f)

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

    return " ".join(text_runs)

uploaded = files.upload()

pptx_files_content = {}
for file_name, file_content in uploaded.items():
    pptx_files_content[file_name] = process_pptx_content(file_content)

df_pptx = pd.DataFrame(pptx_files_content.items(), columns=['file_name', 'text'])

df_pptx_clean = df_pptx[df_pptx['text'].str.len() > 0].reset_index(drop=True)
df_pptx_clean['text_clean'] = df_pptx_clean['text'].apply(lambda x: re.sub(r'[^a-zA-Z0-9\s]+', '', str(x)))

df_pptx_clean.head()

!pip install python-pptx

from docx import Document

from google.colab import files
uploaded = files.upload()

docx_file_name = list(uploaded.keys())[0]

document = Document(docx_file_name)
text = []

for paragraph in document.paragraphs:
    text.append(paragraph.text)

full_text = "\n".join(text)
print(full_text)

!pip install python-docx

!pip install PyMuPDF